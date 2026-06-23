"""
文件解析工具模块
统一接口支持多种文件格式的文本提取
支持: TXT, JSON, JSONL, MD, CSV, PDF, DOCX, DOC, XLSX, XLS, PPTX, RTF, XML, HTML, EPUB, LOG
      INI, YAML, Python/JS/TS/Java/C/C++/Go/Rust/SQL 等源码文件
      PNG, JPG, JPEG, BMP, GIF, WEBP, TIFF（OCR 图片文字识别）
"""
import logging
import os
from typing import Callable, Optional

logger = logging.getLogger("FileParser")

SUPPORTED_EXTENSIONS = {
    ".txt", ".json", ".jsonl", ".md", ".csv", ".pdf", ".docx", ".doc",
    ".html", ".htm", ".epub", ".xlsx", ".xls", ".pptx", ".rtf",
    ".xml", ".log", ".ini", ".cfg", ".yaml", ".yml", ".py", ".js",
    ".ts", ".css", ".java", ".c", ".cpp", ".h", ".hpp", ".sh",
    ".bat", ".ps1", ".sql", ".r", ".go", ".rs", ".swift",
    ".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif"}


def parse_file_to_text(file_path: str,
                       progress_callback: Optional[Callable] = None) -> str:
    """
    解析文件为纯文本

    Args:
        file_path: 文件路径
        progress_callback: 可选进度回调 callback(current, total, message)
                           用于 PDF OCR 场景推送 SSE 进度
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext in {".txt", ".md", ".json", ".jsonl", ".log", ".ini",
               ".cfg", ".yaml", ".yml", ".py", ".js", ".ts", ".css",
               ".java", ".c", ".cpp", ".h", ".hpp", ".sh", ".bat",
               ".ps1", ".sql", ".r", ".go", ".rs", ".swift"}:
        return _parse_text(file_path)
    elif ext == ".csv":
        return _parse_csv(file_path)
    elif ext == ".pdf":
        return _parse_pdf_with_ocr_fallback(file_path, progress_callback)
    elif ext == ".docx":
        return _parse_docx(file_path)
    elif ext == ".doc":
        return _parse_doc(file_path)
    elif ext in {".html", ".htm"}:
        return _parse_html(file_path)
    elif ext == ".epub":
        return _parse_epub(file_path)
    elif ext in {".xlsx", ".xls"}:
        return _parse_excel(file_path)
    elif ext == ".pptx":
        return _parse_pptx(file_path)
    elif ext == ".rtf":
        return _parse_rtf(file_path)
    elif ext == ".xml":
        return _parse_xml(file_path)
    elif ext in IMAGE_EXTENSIONS:
        return _parse_image_ocr(file_path)
    else:
        logger.warning(f"不支持的文件格式: {ext}，尝试作为纯文本读取")
        return _parse_text(file_path)


# ======================== 纯文本 ========================

def _parse_text(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="gbk") as f:
                return f.read()
        except Exception as e:
            logger.error(f"读取文本文件失败: {e}")
            return ""


# ======================== CSV ========================

def _parse_csv(file_path: str) -> str:
    try:
        import pandas as pd
        df = pd.read_csv(file_path)
        lines = []
        for _, row in df.iterrows():
            line = "，".join([f"{col}: {val}" for col, val in row.items()])
            lines.append(line)
        return "\n".join(lines)
    except ImportError:
        return _parse_text(file_path)
    except Exception as e:
        logger.error(f"读取 CSV 失败: {e}")
        return ""


# ======================== PDF（含 OCR 回退） ========================

def _parse_pdf_with_ocr_fallback(file_path: str,
                                 progress_callback: Optional[Callable] = None) -> str:
    basename = os.path.basename(file_path)

    # 第一步：PyMuPDF 文本层
    text = _parse_pdf_pymupdf(file_path).strip()
    if text:
        logger.info(f"PDF (PyMuPDF) 文本提取成功: {basename} ({len(text)} 字符)")
        return text

    # 第二步：PyPDF2 / pdfminer
    text = _parse_pdf_text(file_path).strip()
    if text:
        return text

    # 第三步：OCR（逐页，支持进度回调）
    logger.info(f"PDF 文本层为空，启动 OCR 识别: {basename}")
    ocr_text = _parse_pdf_ocr(file_path, progress_callback=progress_callback)
    if ocr_text.strip():
        return ocr_text

    logger.warning(f"⚠️ 无法从 PDF 提取任何文本: {basename}")
    return ""


def _parse_pdf_pymupdf(file_path: str) -> str:
    try:
        import fitz
        doc = fitz.open(file_path)
        total_pages = len(doc)
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text()
            if text and text.strip():
                pages.append(text)
        doc.close()
        result = "\n\n".join(pages)
        if result:
            logger.info(f"PyMuPDF 文本层: {len(pages)}/{total_pages} 页有文字")
        else:
            logger.info(f"PyMuPDF 文本层: 0/{total_pages} 页 → 需 OCR")
        return result
    except ImportError:
        return ""
    except Exception as e:
        logger.warning(f"PyMuPDF 读取失败: {e}")
        return ""


def _parse_pdf_text(file_path: str) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        pages = [page.extract_text() for page in reader.pages if page.extract_text()]
        return "\n\n".join(pages)
    except (ImportError, Exception):
        return _parse_pdf_fallback(file_path)


def _parse_pdf_fallback(file_path: str) -> str:
    try:
        from pdfminer.high_level import extract_text
        return extract_text(file_path)
    except ImportError:
        return ""


def _parse_pdf_ocr(file_path: str,
                   progress_callback: Optional[Callable] = None) -> str:
    """
    逐页 OCR 识别 PDF（全部页面，支持进度回调）
    使用 ThreadPoolExecutor 并行处理多页以加速。
    添加快速抽样检测：先 OCR 第 1 页，若为空则提前终止。

    Args:
        file_path: PDF 路径
        progress_callback: callback(current_page, total_pages, page_text_or_none)
                           每完成一页调用一次
    """
    try:
        import fitz
        from concurrent.futures import ThreadPoolExecutor, as_completed
        from PIL import Image
        import io
        import threading

        doc = fitz.open(file_path)
        total_pages = len(doc)

        logger.info(f"PDF OCR: 全部 {total_pages} 页（150 DPI，并行 OCR）...")
        if progress_callback:
            progress_callback(0, total_pages, None)

        # ── 快速抽样检测：先渲染第 1 页测试 OCR ──
        if total_pages > 0:
            sample_page = doc[0]
            sample_mat = fitz.Matrix(150 / 72, 150 / 72)
            sample_pix = sample_page.get_pixmap(matrix=sample_mat)
            sample_img_data = sample_pix.tobytes("png")
            sample_img = Image.open(io.BytesIO(sample_img_data))
            sample_text = _ocr_from_pil_image(sample_img)
            if not sample_text.strip():
                doc.close()
                logger.warning(
                    f"PDF OCR 快速抽样检测：第 1 页识别为空，"
                    f"跳过全部 {total_pages} 页（该 PDF 可能为空白/加密/纯图片无文字）"
                )
                if progress_callback:
                    for p in range(1, total_pages + 1):
                        progress_callback(p, total_pages, None)
                return ""

        # ── 并行渲染 + 串行 OCR ──
        # 注意：Tesseract 不是线程安全的！所有 Tesseract 调用必须串行化。
        # 但渲染（fitz get_pixmap）可以并行。
        OCR_DPI = 150  # 稳定清晰，像素量约为 200 DPI 的 56%
        OCR_LOCK = threading.Lock()  # Tesseract 全局锁

        texts = [None] * total_pages
        results_lock = threading.Lock()

        def _render_one_page(page_idx: int) -> tuple:
            """渲染单页为 PIL Image（可并行）"""
            try:
                page = doc[page_idx]
                mat = fitz.Matrix(OCR_DPI / 72, OCR_DPI / 72)
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                return (page_idx, img)
            except Exception as e:
                logger.debug(f"渲染第 {page_idx+1} 页失败: {e}")
                return (page_idx, None)

        # 阶段1：并行渲染所有页面
        MAX_WORKERS = min(4, os.cpu_count() or 2)
        with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            futures = {executor.submit(_render_one_page, i): i for i in range(total_pages)}
            rendered = {}
            for future in as_completed(futures):
                page_idx, img = future.result()
                with results_lock:
                    rendered[page_idx] = img

        # 阶段2：串行 OCR（Tesseract 非线程安全）
        for i in range(total_pages):
            img = rendered.get(i)
            if img is None:
                page_text = ""
            else:
                with OCR_LOCK:  # 串行化 Tesseract 调用
                    page_text = _ocr_from_pil_image(img)
                if not page_text.strip():
                    page_text = ""

            with results_lock:
                texts[i] = page_text if page_text.strip() else ""

            # 回调进度
            if progress_callback:
                progress_callback(i + 1, total_pages,
                                  page_text if page_text.strip() else None)

        doc.close()

        # 按页序组装结果
        result_pages = []
        for i, page_text in enumerate(texts):
            if page_text and page_text.strip():
                result_pages.append(f"【PDF 第 {i+1} 页 OCR】\n{page_text.strip()}")

        if result_pages:
            logger.info(f"PDF OCR 完成: {len(result_pages)}/{total_pages} 页识别到文字")
            return "\n\n".join(result_pages)
        else:
            logger.warning(f"PDF OCR 完成但 0/{total_pages} 页识别到文字")
            return ""

    except ImportError:
        logger.error("PyMuPDF (fitz) 未安装")
        return ""
    except Exception as e:
        logger.error(f"PDF OCR 致命错误: {e}", exc_info=True)
        return ""


# ======================== DOCX / DOC ========================

def _parse_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells_text = [cell.text for cell in row.cells if cell.text.strip()]
                if cells_text:
                    paras.append(" | ".join(cells_text))
        return "\n\n".join(paras)
    except ImportError:
        logger.error("python-docx 未安装")
        return ""
    except Exception as e:
        logger.error(f"读取 docx 失败: {e}")
        return ""


def _parse_doc(file_path: str) -> str:
    try:
        import textract
        text = textract.process(file_path).decode("utf-8", errors="replace")
        if text.strip():
            return text
    except (ImportError, Exception):
        pass
    try:
        import olefile
        ole = olefile.OleFileIO(file_path)
        if ole.exists("WordDocument"):
            for stream_name in ["1Table", "0Table"]:
                if ole.exists(stream_name):
                    raw = ole.openstream(stream_name).read()
                    try:
                        text = raw.decode("utf-16-le", errors="ignore")
                        import re
                        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
                        text = re.sub(r'\s+', ' ', text).strip()
                        if len(text) > 20:
                            ole.close()
                            return text
                    except Exception:
                        pass
            ole.close()
    except (ImportError, Exception):
        pass
    logger.warning(f"无法解析 .doc: {os.path.basename(file_path)}")
    return ""


# ======================== HTML ========================

def _html_to_text_regex(html: str) -> str:
    """HTML 转纯文本（纯 regex，无外部依赖）"""
    import re
    text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<nav[^>]*>.*?</nav>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<footer[^>]*>.*?</footer>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<header[^>]*>.*?</header>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<aside[^>]*>.*?</aside>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<(?:br|p|div|h[1-6]|li|tr|blockquote)[^>]*/?>', '\n', text, flags=re.IGNORECASE)
    text = re.sub(r'</(?:p|div|h[1-6]|li|tr|blockquote)>', '\n', text, flags=re.IGNORECASE)
    text = re.sub(r'<[^>]+>', '', text)
    text = text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>').replace('&quot;', '"').replace('&#39;', "'").replace('&nbsp;', ' ')
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _parse_html(file_path: str) -> str:
    try:
        import re
        with open(file_path, "r", encoding="utf-8") as f:
            html = f.read()
        return re.sub(r'\n{3,}', '\n\n', _html_to_text_regex(html))
    except Exception as e:
        logger.error(f"读取 HTML 失败: {e}")
        return _parse_text(file_path)


# ======================== EPUB ========================

def _parse_epub(file_path: str) -> str:
    try:
        from ebooklib import epub
        book = epub.read_epub(file_path)
        chapters = []
        for item in book.get_items_of_type(9):
            html = item.get_content().decode('utf-8', errors='ignore')
            text = _html_to_text_regex(html)
            if text:
                chapters.append(text)
        return "\n\n".join(chapters)
    except ImportError:
        return _parse_epub_fallback(file_path)
    except Exception as e:
        logger.error(f"读取 EPUB 失败: {e}")
        return _parse_epub_fallback(file_path)


def _parse_epub_fallback(file_path: str) -> str:
    import zipfile, re
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            html_files = [n for n in zf.namelist()
                          if n.endswith(('.html', '.xhtml', '.htm'))
                          and not n.startswith('__MACOSX')]
            chapters = []
            for name in sorted(html_files):
                try:
                    content = zf.read(name).decode("utf-8")
                    text = _html_to_text_regex(content)
                    if text:
                        chapters.append(re.sub(r'\n{3,}', '\n\n', text))
                except Exception:
                    continue
            return "\n\n".join(chapters)
    except Exception as e:
        logger.error(f"EPUB 备用解析失败: {e}")
        return ""


# ======================== Excel ========================

def _parse_excel(file_path: str) -> str:
    try:
        import pandas as pd
        xl = pd.ExcelFile(file_path)
        all_sheets = []
        for sheet_name in xl.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            if df.empty:
                continue
            all_sheets.append(f"【工作表: {sheet_name}】")
            all_sheets.append(" | ".join([str(c) for c in df.columns]))
            for _, row in df.head(500).iterrows():
                line = " | ".join([str(v) for v in row.values if pd.notna(v)])
                if line.strip():
                    all_sheets.append(line)
        return "\n".join(all_sheets)
    except ImportError:
        return _parse_excel_openpyxl(file_path)
    except Exception as e:
        logger.error(f"pandas 读取 Excel 失败: {e}")
        return _parse_excel_openpyxl(file_path)


def _parse_excel_openpyxl(file_path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        all_sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_sheets.append(f"【工作表: {sheet_name}】")
            for row in ws.iter_rows(values_only=True):
                line = " | ".join([str(c) for c in row if c is not None])
                if line.strip():
                    all_sheets.append(line)
        wb.close()
        return "\n".join(all_sheets)
    except (ImportError, Exception) as e:
        logger.error(f"openpyxl 失败: {e}")
        return ""


# ======================== PPTX ========================

def _parse_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))
            if slide_texts:
                slides.append(f"【幻灯片 {i}】\n" + "\n".join(slide_texts))
        return "\n\n".join(slides)
    except (ImportError, Exception) as e:
        logger.error(f"读取 pptx 失败: {e}")
        return ""


# ======================== RTF ========================

def _parse_rtf(file_path: str) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        with open(file_path, "r", encoding="utf-8") as f:
            return rtf_to_text(f.read())
    except ImportError:
        return _parse_rtf_fallback(file_path)
    except Exception as e:
        logger.error(f"读取 RTF 失败: {e}")
        return _parse_rtf_fallback(file_path)


def _parse_rtf_fallback(file_path: str) -> str:
    import re
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        content = re.sub(r'\\[a-zA-Z]+\d* ?', '', content)
        content = re.sub(r'\\\'[0-9a-fA-F]{2}', '', content)
        content = re.sub(r'[{}]', '', content)
        return re.sub(r'\n{3,}', '\n\n', content).strip()
    except Exception as e:
        logger.error(f"RTF 备用解析失败: {e}")
        return ""


# ======================== XML ========================

def _parse_xml(file_path: str) -> str:
    try:
        import xml.etree.ElementTree as ET
        tree = ET.parse(file_path)
        root = tree.getroot()

        def _extract_text(elem, depth=0):
            texts = []
            if elem.text and elem.text.strip():
                tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                texts.append(f"{'  ' * depth}[{tag}] {elem.text.strip()}")
            for child in elem:
                texts.extend(_extract_text(child, depth + 1))
            return texts

        lines = _extract_text(root)
        if not lines:
            import re
            with open(file_path, "r", encoding="utf-8") as f:
                xml_content = f.read()
            text = re.sub(r'<[^>]+>', ' ', xml_content)
            text = re.sub(r'\s+', ' ', text).strip()
            return text
        return "\n".join(lines)
    except (ImportError, Exception):
        return _parse_text(file_path)


# ======================== 图片 OCR ========================

def _parse_image_ocr(file_path: str) -> str:
    try:
        from PIL import Image
        img = Image.open(file_path)
        return _ocr_from_pil_image(img, file_path=file_path)
    except (ImportError, Exception) as e:
        logger.error(f"打开图片失败 {file_path}: {e}")
        return ""


def _ocr_from_pil_image(img, file_path: str = None) -> str:
    """从 PIL Image OCR 识别，含预处理增强扫描件准确率

    OCR 回退链：
    1. Tesseract（subprocess 直接调用，带 timeout，避免 pytesseract 死锁）
    2. Windows 原生 OCR API（兜底，Win10+ 零依赖）

    成功时返回识别文本；所有引擎不可用时返回空字符串，
    并调用 _log_ocr_diagnostic() 输出完整安装指引到日志。
    """
    _auto_detect_tesseract()

    tesseract_available = False
    tesseract_path = None
    try:
        import pytesseract
        current = pytesseract.pytesseract.tesseract_cmd
        if current and os.path.isfile(current):
            tesseract_available = True
            tesseract_path = current
    except Exception:
        pass

    # ── 第一优先：Tesseract OCR（subprocess 直接调用，防死锁）──
    if tesseract_available:
        try:
            import subprocess, tempfile
            # 保存图片到临时文件（避免管道死锁）
            _fd, _tmp_img = tempfile.mkstemp(suffix=".png")
            os.close(_fd)
            img.save(_tmp_img, "PNG")

            _fd2, _tmp_out = tempfile.mkstemp(suffix=".txt")
            os.close(_fd2)

            tesseract_exe = os.path.abspath(tesseract_path)
            # 直接调用 tesseract.exe，30 秒超时，隐藏黑窗口
            proc = subprocess.run(
                [tesseract_exe, _tmp_img, _tmp_out.replace(".txt", ""), "-l", "chi_sim+eng"],
                capture_output=True, timeout=30, text=True,
                creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0,
            )
            result_path = _tmp_out if _tmp_out.endswith(".txt") else _tmp_out + ".txt"
            text = ""
            if os.path.exists(result_path):
                with open(result_path, "r", encoding="utf-8") as f:
                    text = f.read().strip()

            # 清理临时文件
            try: os.remove(_tmp_img)
            except: pass
            try: os.remove(result_path)
            except: pass

            if text:
                return text

        except subprocess.TimeoutExpired:
            logger.debug(f"Tesseract 超时（30s）")
        except Exception as e:
            logger.debug(f"Tesseract subprocess 调用失败: {e}")
    else:
        logger.warning(
            "Tesseract-OCR 未安装或路径未检测到。"
            "请从 https://github.com/UB-Mannheim/tesseract/wiki 下载安装，"
            "并勾选中文语言包 chi_sim。"
        )

    # ── 第二优先：Windows 原生 OCR（Win10+ 零依赖兜底）──
    windows_ocr_available = False
    try:
        import winrt
        windows_ocr_available = True
    except ImportError:
        logger.warning("winrt 库未安装，Windows 原生 OCR 不可用。pip install winrt")

    _temp_path = None
    if windows_ocr_available:
        try:
            if not file_path:
                # PDF 渲染的图片没有实体文件路径，临时写入
                import tempfile
                _temp_fd, _temp_path = tempfile.mkstemp(suffix=".png")
                os.close(_temp_fd)
                img.save(_temp_path, "PNG")
                file_path = _temp_path

            win_text = _windows_native_ocr(file_path)
            if win_text.strip():
                logger.info("Windows 原生 OCR 兜底成功")
                return win_text
        except Exception:
            pass
        finally:
            if _temp_path:
                try:
                    os.remove(_temp_path)
                except Exception:
                    pass

    # ── 所有引擎尝试完毕，仍无文字 ──
    if not tesseract_available and not windows_ocr_available:
        # 引擎完全缺失：输出安装指引
        _log_ocr_diagnostic(tesseract_available, tesseract_path, windows_ocr_available)
    else:
        # 引擎可用但识别为空：该页可能为空白/纯图/低质量扫描
        logger.warning(
            "OCR 引擎已就绪但未识别到文字，该页面可能为空白、纯图片或质量过低的扫描件。"
        )
    return ""


def _probe_ocr_engines() -> dict:
    """探测所有 OCR 引擎可用状态，返回字典供外部（SSE/API）展示诊断信息"""
    import pytesseract

    _auto_detect_tesseract()

    tesseract_available = False
    tesseract_path = None
    tesseract_langs = []
    try:
        current = pytesseract.pytesseract.tesseract_cmd
        if current and os.path.isfile(current):
            tesseract_available = True
            tesseract_path = current
            try:
                tesseract_langs = pytesseract.get_languages()
            except Exception:
                tesseract_langs = ["eng"]  # 无法探测时假设至少支持英文
    except Exception:
        pass

    has_chi_sim = "chi_sim" in tesseract_langs

    windows_ocr_available = False
    try:
        import winrt
        windows_ocr_available = True
    except ImportError:
        pass

    return {
        "tesseract_available": tesseract_available,
        "tesseract_path": tesseract_path,
        "tesseract_langs": tesseract_langs,
        "has_chi_sim": has_chi_sim,
        "windows_ocr_available": windows_ocr_available,
    }


def get_ocr_diagnostic_text() -> str:
    """
    返回 OCR 引擎状态的纯文本诊断信息（用于 SSE 推送到前端）。
    无论 OCR 是否可用，始终返回可读状态。
    """
    status = _probe_ocr_engines()
    lines = [
        "🔍 OCR 引擎状态检测",
        "",
    ]

    if status["tesseract_available"]:
        lines.append(f"  ✅ Tesseract-OCR: {status['tesseract_path']}")
        if status["has_chi_sim"]:
            lines.append(f"     ✅ 中文语言包 chi_sim 已安装")
        else:
            lines.append(f"     ❌ 缺少中文语言包 chi_sim！可用语言: {status['tesseract_langs']}")
            lines.append(f"     请下载 chi_sim.traineddata 放入 Tesseract tessdata 目录")
    else:
        lines.append(f"  ❌ Tesseract-OCR 未安装")
        lines.append(f"     → 下载: https://github.com/UB-Mannheim/tesseract/wiki")
        lines.append(f"     安装时必须勾选 Chinese (Simplified) 语言包")

    if status["windows_ocr_available"]:
        lines.append(f"  ✅ Windows 原生 OCR (winrt) 已安装")
    else:
        lines.append(f"  ❌ Windows 原生 OCR (winrt) 未安装")
        lines.append(f"     → 命令: pip install winrt")

    if not status["tesseract_available"] and not status["windows_ocr_available"]:
        lines.append("")
        lines.append("  ⚠️ 两个 OCR 引擎均不可用，扫描版 PDF / 图片将无法提取文字！")
        lines.append("  请按上述指引安装至少一个 OCR 引擎后重启应用。")
    else:
        # 引擎可用但仍可能 0 字符：可能是空白页/低质量扫描
        lines.append("")
        lines.append("  ℹ️ OCR 引擎已就绪。如果 PDF 仍未提取到文字，可能是：")
        lines.append("     1. 该 PDF 页面为空白页")
        lines.append("     2. 扫描质量过低（模糊/过暗）")
        lines.append("     3. 该 PDF 为加密/受保护文件")

    return "\n".join(lines)


def _log_ocr_diagnostic(tesseract_available: bool, tesseract_path: str | None,
                        windows_ocr_available: bool) -> None:
    """输出 OCR 引擎不可用的详细诊断信息和安装指引"""
    lines = [
        "=" * 60,
        "⚠️  OCR 引擎不可用 — 扫描版 PDF/图片无法识别文字",
        "=" * 60,
        f"  Tesseract-OCR: {'✅ ' + (tesseract_path or '') if tesseract_available else '❌ 未检测到'}",
        f"  Windows OCR:   {'✅ 已安装' if windows_ocr_available else '❌ winrt 未安装'}",
        "",
        "  请选择以下方案之一配置 OCR 引擎：",
        "",
        "  方案一 [推荐] 安装 Tesseract-OCR：",
        "    → https://github.com/UB-Mannheim/tesseract/wiki",
        "    安装时必须勾选 Chinese (Simplified) 语言包",
        "",
        "  方案二 安装 Windows 原生 OCR（零额外依赖）：",
        "    → pip install winrt",
        "",
        "  配置完成后重启本应用即可。",
        "=" * 60,
    ]
    for line in lines:
        logger.warning(line)


def _windows_native_ocr(image_path: str) -> str:
    """
    Windows 10/11 原生 OCR API 兜底方案。
    通过 Windows.Media.Ocr 引擎识别图片文字，零额外依赖，零体积。
    支持中英文混合识别。
    返回空字符串表示不可用或失败。
    """
    import sys as _sys
    if _sys.platform != "win32":
        return ""
    try:
        import asyncio
        from winrt.windows.media.ocr import OcrEngine
        from winrt.windows.graphics.imaging import (
            BitmapDecoder, SoftwareBitmap, BitmapPixelFormat, BitmapAlphaMode
        )
        from winrt.windows.storage import StorageFile

        async def _ocr():
            abs_path = os.path.abspath(image_path)
            try:
                file = await StorageFile.get_file_from_path_async(abs_path)
            except Exception:
                return ""

            stream = await file.open_async(0)  # FileAccessMode.Read
            decoder = await BitmapDecoder.create_async(stream)

            # 尝试获取适合尺寸的 bitmap
            try:
                bitmap = await decoder.get_software_bitmap_async()
            except Exception:
                # fallback：有些格式可能需要转换像素格式
                bitmap = await decoder.get_software_bitmap_async(
                    BitmapPixelFormat.bgra8,
                    BitmapAlphaMode.premultiplied
                )

            engine = OcrEngine.try_create_from_user_profile_languages()
            if engine is None:
                from winrt.windows.globalization import Language
                try:
                    engine = OcrEngine.try_create_from_language(Language("zh-Hans"))
                except Exception:
                    pass
            if engine is None:
                try:
                    engine = OcrEngine.try_create_from_language(Language("en-US"))
                except Exception:
                    pass
            if engine is None:
                logger.debug("Windows OCR 引擎不可用（无可识别语言）")
                return ""

            result = await engine.recognize_async(bitmap)
            lines = [line.text for line in result.lines]
            return "\n".join(lines)

        return asyncio.run(_ocr())

    except ImportError:
        logger.debug("winrt 库未安装，Windows 原生 OCR 不可用")
        return ""
    except Exception as e:
        logger.debug(f"Windows 原生 OCR 失败: {e}")
        return ""


def _auto_detect_tesseract():
    import os as _os, sys
    if sys.platform != "win32":
        return
    try:
        import pytesseract
    except ImportError:
        return
    try:
        current = pytesseract.pytesseract.tesseract_cmd
        if current and _os.path.exists(current):
            return
    except Exception:
        pass
    candidates = [
        _os.path.join(_os.environ.get("ProgramFiles", "C:\\Program Files"), "Tesseract-OCR", "tesseract.exe"),
        _os.path.join(_os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)"), "Tesseract-OCR", "tesseract.exe"),
        _os.path.join(_os.environ.get("LOCALAPPDATA", ""), "Programs", "Tesseract-OCR", "tesseract.exe"),
        "C:\\Program Files\\Tesseract-OCR\\tesseract.exe",
        "C:\\Program Files (x86)\\Tesseract-OCR\\tesseract.exe",
    ]
    for c in candidates:
        if _os.path.isfile(c):
            pytesseract.pytesseract.tesseract_cmd = c
            logger.info(f"Tesseract: {c}")
            return
